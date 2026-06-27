"""Text extraction for AI summarization.

Course materials are mostly PDF / Office docs whose raw bytes are not UTF-8 text,
so feeding them straight to the model summarizes binary noise. This pulls readable
text per file type before the AI call. Parsers are lazy-imported and any failure
degrades to a best-effort UTF-8 decode, so a missing dep or a malformed file never
crashes the summary endpoint — it just yields whatever text it can.
"""

from __future__ import annotations

import io

# Cap so we never hand the model a whole textbook; summarize() truncates again.
MAX_CHARS = 20000


def _from_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    parts: list[str] = []
    for page in reader.pages:
        parts.append(page.extract_text() or "")
        if sum(len(p) for p in parts) >= MAX_CHARS:
            break
    return "\n".join(parts)


def _from_docx(data: bytes) -> str:
    import docx  # python-docx

    doc = docx.Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs)


def _from_pptx(data: bytes) -> str:
    from pptx import Presentation  # python-pptx

    prs = Presentation(io.BytesIO(data))
    lines: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                lines.append(shape.text_frame.text)
    return "\n".join(lines)


# Suffix -> extractor. Plain-text types fall through to the UTF-8 decode below.
_PARSERS = {
    ".pdf": _from_pdf,
    ".docx": _from_docx,
    ".pptx": _from_pptx,
}


def extract_text(filename: str, data: bytes) -> str:
    """Return readable text for `data`, chosen by `filename`'s extension.

    PDF/DOCX/PPTX are parsed; everything else (txt/md/csv/code) is decoded as
    UTF-8. Any parser error falls back to the UTF-8 decode so we still hand the
    model *something* rather than failing the request.
    """
    ext = ""
    if "." in filename:
        ext = "." + filename.rsplit(".", 1)[-1].lower()

    parser = _PARSERS.get(ext)
    if parser is not None:
        try:
            text = parser(data)
            if text.strip():
                return text[:MAX_CHARS]
        except Exception:
            pass  # fall through to raw decode

    return data.decode("utf-8", errors="ignore")[:MAX_CHARS]
